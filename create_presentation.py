from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_deck():
    prs = Presentation()
    
    # We will use the blank layout (usually index 6)
    blank_slide_layout = prs.slide_layouts[6]
    
    # Theme Colors
    BG_COLOR = RGBColor(11, 23, 41)       # #0B1729 (Dark Blue)
    ACCENT_BLUE = RGBColor(31, 111, 235)  # #1F6FEB
    ACCENT_CYAN = RGBColor(0, 200, 255)   # #00C8FF
    TEXT_WHITE = RGBColor(255, 255, 255)  # #FFFFFF
    TEXT_GRAY = RGBColor(200, 200, 200)   # #C8C8C8
    
    def apply_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
    def add_title(slide, text):
        left = Inches(0.5)
        top = Inches(0.4)
        width = Inches(9)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = ACCENT_CYAN
        
        # Add a subtle line under title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(3), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_BLUE
        line.line.color.rgb = ACCENT_BLUE

    def add_speaker_notes(slide, text):
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = text

    # ==========================================
    # SLIDE 1: Title & Vision
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide1)
    
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    txBox = slide1.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Flipkart GridLock 2.0"
    p.font.bold = True
    p.font.size = Pt(54)
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "AI-Powered Traffic Diversion & Incident Management Platform"
    p2.font.size = Pt(24)
    p2.font.color.rgb = ACCENT_CYAN
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "\nPredictive Intelligence for Smart City Resilience"
    p3.font.size = Pt(18)
    p3.font.color.rgb = TEXT_GRAY
    p3.alignment = PP_ALIGN.CENTER
    
    add_speaker_notes(slide1, "Welcome to our presentation on Flipkart GridLock 2.0. This platform is a comprehensive AI-powered traffic diversion and incident management system designed for smart cities. We've built an end-to-end command center that transforms historical data and live inputs into predictive, actionable traffic management decisions.")

    # ==========================================
    # SLIDE 2: Problem & Opportunity
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide2)
    add_title(slide2, "The Problem: Event-Driven Congestion")
    
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(4.2)
    height = Inches(4.5)
    
    # Problem Box 1
    shape1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(4), Inches(2.2))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(20, 40, 70)
    shape1.line.color.rgb = RGBColor(231, 76, 60) # Red accent
    tf1 = shape1.text_frame
    p = tf1.paragraphs[0]
    p.text = "Traffic Congestion Challenges"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_WHITE
    p = tf1.add_paragraph()
    p.text = "• Chronic blackspots & severe bottlenecks\n• Reactive instead of proactive management\n• 4,800+ breakdowns causing cascaded delays"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_GRAY

    # Problem Box 2
    shape2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.3), Inches(4), Inches(2.2))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(20, 40, 70)
    shape2.line.color.rgb = RGBColor(231, 76, 60)
    tf2 = shape2.text_frame
    p = tf2.paragraphs[0]
    p.text = "Incident Response Delays"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_WHITE
    p = tf2.add_paragraph()
    p.text = "• Manual, unoptimized manpower allocation\n• Lack of real-time multi-factor modeling\n• Over or under-deployment of police & barricades"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_GRAY

    # Problem Box 3
    shape3 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(1.8), Inches(4.5), Inches(4.7))
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = RGBColor(20, 40, 70)
    shape3.line.color.rgb = RGBColor(241, 196, 15) # Yellow accent
    tf3 = shape3.text_frame
    p = tf3.paragraphs[0]
    p.text = "The Opportunity"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_CYAN
    p = tf3.add_paragraph()
    p.text = "\nBy leveraging AI and historical Astram data (8,200+ incidents), we can predict clearance times, mine chronic blackspots, and automate diversion planning. Moving from reactive crisis control to predictive smart-city management saves fuel, time, and reduces emissions."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_GRAY
    
    add_speaker_notes(slide2, "Bengaluru faces massive event-driven congestion. Currently, traffic management is reactive, leading to severe bottlenecks and delayed incident response. By moving to a predictive model using our historical Astram data, we have the opportunity to proactively manage traffic, allocate resources intelligently, and execute dynamic diversions.")

    # ==========================================
    # SLIDE 3: Our Solution
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide3)
    add_title(slide3, "Our Solution: End-to-End Workflow")
    
    # Workflow Diagram
    steps = [
        ("1. Data Fusion", "Ingest Astram data, weather, & live events"),
        ("2. AI Prediction", "XGBoost & LightGBM predict clearance & severity"),
        ("3. Optimization", "OR-Tools allocates manpower & barricades"),
        ("4. Diversion", "NetworkX plans optimal offline detours"),
        ("5. Command Center", "Streamlit visualizes actionable intelligence")
    ]
    
    for i, (title, desc) in enumerate(steps):
        left = Inches(0.5 + i * 1.8)
        top = Inches(3.0)
        shape = slide3.shapes.add_shape(MSO_SHAPE.HEXAGON, left, top, Inches(1.6), Inches(1.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(31, 111, 235)
        shape.line.color.rgb = ACCENT_CYAN
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_WHITE
        p.alignment = PP_ALIGN.CENTER
        
        top_desc = Inches(4.8)
        txBox = slide3.shapes.add_textbox(left, top_desc, Inches(1.6), Inches(1.5))
        tf_desc = txBox.text_frame
        tf_desc.word_wrap = True
        p_desc = tf_desc.paragraphs[0]
        p_desc.text = desc
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_GRAY
        p_desc.alignment = PP_ALIGN.CENTER
        
        if i < len(steps) - 1:
            arrow = slide3.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.1 + i * 1.8), Inches(3.6), Inches(0.3), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_CYAN
            arrow.line.color.rgb = ACCENT_CYAN
            
    add_speaker_notes(slide3, "Our platform provides an end-to-end workflow. We fuse real-time event parameters with historical Astram data and weather APIs. Our AI models predict clearance time and priority. We then use OR-Tools to optimize police dispatch, compute diversions using a spatial network graph, and display everything in a highly interactive Command Center.")

    # ==========================================
    # SLIDE 4: AI Intelligence Engine
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide4)
    add_title(slide4, "AI Intelligence Engine")
    
    # 4 quadrants of AI
    quads = [
        (Inches(0.5), Inches(2.0), "XGBoost Congestion Model", "Predicts congestion score (0-100) using event type, location, weather, crowd size, and diurnal curves. R² = 0.889"),
        (Inches(5.0), Inches(2.0), "LightGBM Clearance Regressor", "Estimates exact incident clearance time (MAE = 74 min). Fuses text fields via SentenceTransformers (NLP)."),
        (Inches(0.5), Inches(4.5), "Surge & Blackspot Mining", "DBSCAN clusters chronic incident hotspots. Anomaly detection algorithms identify volume surges in real-time."),
        (Inches(5.0), Inches(4.5), "Explainable AI (XAI)", "Provides transparent reasoning chains for decisions, showing feature weights (e.g., 'Historical Baseline' influenced 50%).")
    ]
    
    for left, top, title, desc in quads:
        shape = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(4.3), Inches(1.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(15, 30, 55)
        shape.line.color.rgb = ACCENT_BLUE
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = ACCENT_CYAN
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_GRAY
        
    add_speaker_notes(slide4, "The brain of GridLock 2.0 uses powerful ML models. Our XGBoost engine predicts congestion with 89% R-squared. We use LightGBM for clearance time prediction, enriched with NLP embeddings. DBSCAN helps us map blackspots. Most importantly, our Explainable AI (XAI) transparently outlines exactly why it made a deployment recommendation.")

    # ==========================================
    # SLIDE 5: Technical Architecture
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide5)
    add_title(slide5, "Technical Architecture")
    
    layers = [
        ("Presentation Layer", "Streamlit Dashboard • PyDeck 3D Maps • Plotly Charts", Inches(1.5)),
        ("Service Layer", "FastAPI Endpoints • Traffic Engine • KPI Engine • Spatial Engine", Inches(3.0)),
        ("AI / ML Layer", "XGBoost • LightGBM • OR-Tools (CP-SAT) • NetworkX Routing", Inches(4.5)),
        ("Data Layer", "SQLite DB • Pandas • Cleaned Astram CSVs • Open-Meteo", Inches(6.0))
    ]
    
    for title, desc, top in layers:
        shape = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), top, Inches(8), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(20, 40, 70)
        shape.line.color.rgb = ACCENT_CYAN
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_WHITE
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_GRAY
        p2.alignment = PP_ALIGN.CENTER
        
    add_speaker_notes(slide5, "Our architecture is highly modular and strictly open-source. The Presentation Layer uses Streamlit for a fast, beautiful command center. The Service Layer is powered by FastAPI and our custom Traffic Engine. The ML layer runs XGBoost and Google OR-Tools. Data is managed through SQLite and Pandas, ensuring the app can run fully offline if needed.")

    # ==========================================
    # SLIDE 6: Traffic Command Center
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide6)
    add_title(slide6, "Traffic Command Center UI")
    
    # We describe the 7 screens instead of screenshots (since we don't have images)
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9)
    height = Inches(4.5)
    txBox = slide6.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    screens = [
        "1. City Overview: Live PyDeck 3D map plotting incidents, diversions, and deployments.",
        "2. Impact Prediction: Real-time congestion gauge and expected delay metrics.",
        "3. Live Heatmap: Multi-factor traffic simulation + 15/30/60 min AR(1) forecasting.",
        "4. Diversion Management: Alternate routes computed via NetworkX with fuel & time savings.",
        "5. Fleet Dispatch: Gap analysis and OR-Tools optimization for officers, marshals, and barricades.",
        "6. AI Commander: Executive narrative explaining feature weights and system confidence.",
        "7. Retraining Loop: Post-event accuracy report cards and automated SQLite-based learning."
    ]
    
    for s in screens:
        p = tf.add_paragraph()
        p.text = s
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_GRAY
        p.space_after = Pt(10)
        
    add_speaker_notes(slide6, "Our dashboard is a comprehensive 7-screen Command Center. It offers a 3D city overview, real-time congestion gauges, a live heatmap with AR(1) forecasting, dedicated diversion management, fleet dispatch maps, an AI Commander for explainability, and a continuous learning loop module.")

    # ==========================================
    # SLIDE 7: Key Innovations
    # ==========================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide7)
    add_title(slide7, "Key Innovations")
    
    innovations = [
        ("Dynamic Traffic Engine", "Fuses diurnal curves, weather, event types, and closures to simulate realistic minute-by-minute traffic."),
        ("Micro-Simulation Pipeline", "Time-stepped queue modeling estimates emissions & delay, exporting scenarios to SUMO."),
        ("Offline Proximity Graph", "Generates network paths locally from incident GPS, eliminating dependency on constant internet APIs."),
        ("Constraint Optimization", "OR-Tools dynamically balances limited police personnel across high-risk corridors to maximize impact.")
    ]
    
    for i, (title, desc) in enumerate(innovations):
        top = Inches(2.0 + i * 1.3)
        left = Inches(0.5)
        shape = slide7.shapes.add_shape(MSO_SHAPE.RIGHT_BRACE, left, top, Inches(0.3), Inches(1.0))
        shape.line.color.rgb = ACCENT_CYAN
        
        txBox = slide7.shapes.add_textbox(Inches(1.0), top - Inches(0.1), Inches(8), Inches(1.0))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = ACCENT_CYAN
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(16)
        p2.font.color.rgb = TEXT_GRAY
        
    add_speaker_notes(slide7, "What makes GridLock 2.0 innovative? First, our dynamic traffic engine simulates multi-factor live traffic. Second, our custom micro-simulator generates SUMO-ready bundles. Third, we build our own offline proximity graphs so the system never breaks without internet. Finally, we use military-grade constraint optimization to deploy police forces.")

    # ==========================================
    # SLIDE 8: Demo & Results
    # ==========================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide8)
    add_title(slide8, "Demonstrated Performance")
    
    # 3 big metric boxes
    metrics = [
        ("XGBoost R² Score", "0.889", "Highly accurate congestion scoring on synthetic test set", Inches(0.5)),
        ("Model Confidence", ">85%", "For high-risk critical events with explainable feature weights", Inches(3.5)),
        ("Clearance MAE", "74 Min", "Mean Absolute Error on predicting real incident clearance", Inches(6.5))
    ]
    
    for title, val, desc, left in metrics:
        shape = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.5), Inches(2.8), Inches(2.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(15, 30, 55)
        shape.line.color.rgb = ACCENT_BLUE
        
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = val
        p.font.bold = True
        p.font.size = Pt(40)
        p.font.color.rgb = RGBColor(46, 204, 113) # Green
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.bold = True
        p2.font.size = Pt(16)
        p2.font.color.rgb = TEXT_WHITE
        p2.alignment = PP_ALIGN.CENTER
        
        p3 = tf.add_paragraph()
        p3.text = "\n" + desc
        p3.font.size = Pt(12)
        p3.font.color.rgb = TEXT_GRAY
        p3.alignment = PP_ALIGN.CENTER
        
    add_speaker_notes(slide8, "Our system is highly performant. On our benchmark sets, the XGBoost engine achieved an R-squared of 0.889. The system routinely maintains over 85% confidence in deployment recommendations. And for clearance times, our mean absolute error is 74 minutes across thousands of varied Astram records.")

    # ==========================================
    # SLIDE 9: Business Impact
    # ==========================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide9)
    add_title(slide9, "Business & Civic Impact")
    
    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(9)
    height = Inches(4.5)
    txBox = slide9.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    points = [
        "🛣️ Faster Incident Clearance: Intelligent dispatch reduces response times, unblocking arterials sooner.",
        "🌿 Emissions Reduction: Every diverted vehicle saves idling fuel and reduces overall city CO2 footprint.",
        "📊 Efficient Resource Usage: Prevents under-deployment at crisis zones and over-deployment at low-risk events.",
        "📈 Predictive Governance: Allows authorities to transition from crisis management to predictive planning.",
        "🚀 Scalability: Zero-cost open-source stack (SQLite, Pandas, LightGBM) scales instantly to any tier-1 city."
    ]
    
    for s in points:
        p = tf.add_paragraph()
        p.text = s
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_GRAY
        p.space_after = Pt(18)
        
    add_speaker_notes(slide9, "The impact is massive. We clear incidents faster, saving thousands of liters of fuel and CO2 emissions. We optimize exactly how many police and barricades are deployed. Most importantly, it scales. Built on an entirely open-source stack, this platform can be deployed in any tier-1 city worldwide without expensive licensing.")

    # ==========================================
    # SLIDE 10: Future Vision & Closing
    # ==========================================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide10)
    add_title(slide10, "Future Vision & Roadmap")
    
    shape = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.0), Inches(8.5), Inches(3.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(15, 30, 55)
    shape.line.color.rgb = ACCENT_CYAN
    tf = shape.text_frame
    
    future_points = [
        "• V2X Integration: Direct API linking to connected vehicles for automated dynamic rerouting.",
        "• Computer Vision: Ingesting live CCTV feeds to automatically validate crowd sizes & incident severity.",
        "• Dynamic Signal Control: Integrating with smart traffic lights to flush diversion corridors.",
        "• Flipkart Logistics: Re-routing delivery fleets dynamically based on predictive gridlock APIs."
    ]
    
    for s in future_points:
        p = tf.add_paragraph()
        p.text = s
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_GRAY
        p.space_after = Pt(10)
        
    txBox2 = slide10.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(8.5), Inches(1))
    tf2 = txBox2.text_frame
    p_end = tf2.paragraphs[0]
    p_end.text = "Thank You. Ready for deployment."
    p_end.font.bold = True
    p_end.font.size = Pt(32)
    p_end.font.color.rgb = ACCENT_BLUE
    p_end.alignment = PP_ALIGN.CENTER
    
    add_speaker_notes(slide10, "Looking ahead, we envision V2X integration, live CCTV ingestion, and linking this engine to Flipkart's logistics fleet to guarantee delivery ETAs. GridLock 2.0 is more than a dashboard; it's a foundation for smart city traffic governance. Thank you for your time.")

    # Save presentation
    prs.save("Project-Presentation.pptx")

if __name__ == "__main__":
    create_deck()
